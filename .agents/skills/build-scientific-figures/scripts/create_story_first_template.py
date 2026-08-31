#!/usr/bin/env python3
"""Create a neutral editable story-first hero template from the bundled case study."""

from __future__ import annotations

import argparse
from pathlib import Path
import tempfile
from typing import Any, Sequence

from PIL import Image, ImageDraw
from pptx import Presentation


TEXT = {
    "exact-request-value": '"fixed value"',
    "open-request-value": '"allowed range"',
    "exact-intent-text": "v*",
    "open-intent-text": "S",
    "exact-outcome-provenance": "source_A",
    "clarify-outcome-provenance": "NOT ISSUED",
    "declared-outcome-provenance": "source_B",
    "shared-contract-value": "value = v*",
    "execution-metric": "recorded endpoint witness",
    "execution-label-1": "1  State A",
    "execution-label-2": "2  State B",
    "execution-label-3": "3  State C",
    "execution-label-4": "4  State D",
}


def _replace_picture(slide: Any, picture: Any, path: Path) -> None:
    _, relationship_id = slide.part.get_or_add_image_part(str(path))
    picture._pic.blipFill.blip.rEmbed = relationship_id


def _placeholder(path: Path, index: int) -> None:
    image = Image.new("RGB", (960, 720), "#F1F5F9")
    draw = ImageDraw.Draw(image)
    colors = ("#135FD9", "#64748B", "#159F91", "#E17A00")
    color = colors[index - 1]
    x = 210 + (index - 1) * 70
    draw.rounded_rectangle((x, 225, x + 420, 495), radius=42, fill="white", outline=color, width=14)
    draw.ellipse((x + 135, 300, x + 285, 450), fill=color)
    image.save(path, optimize=True)


def build(source: Path, output: Path) -> Path:
    prs = Presentation(source)
    if len(prs.slides) != 1:
        raise ValueError("story-first template source must have one slide")
    slide = prs.slides[0]
    shapes = {shape.name: shape for shape in slide.shapes}
    for name, value in TEXT.items():
        shape = shapes[name]
        shape.text_frame.paragraphs[0].runs[0].text = value
    with tempfile.TemporaryDirectory(prefix="story-first-template-") as temporary:
        root = Path(temporary)
        for index in range(1, 5):
            path = root / f"state-{index}.png"
            _placeholder(path, index)
            _replace_picture(slide, shapes[f"execution-image-{index}"], path)
        output.parent.mkdir(parents=True, exist_ok=True)
        prs.core_properties.title = "Story-first scientific hero template"
        prs.core_properties.subject = "Neutral editable branch-merge-execution figure"
        prs.core_properties.author = "Scientific Figure Builder skill"
        prs.save(output)
    return output


def main(argv: Sequence[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    args = parser.parse_args(argv)
    print(build(args.input.resolve(), args.output.resolve()))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
